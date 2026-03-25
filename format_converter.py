import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox
import os
import threading
import shutil
import pandas as pd
from PIL import Image
import PyPDF2
from pdf2docx import Converter
from docx2pdf import convert
import docx
from pptx import Presentation

try:
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    COMTYPES_AVAILABLE = False

import config
import ai_handler

def list_supported_formats(file_ext):
    ext = file_ext.lower().replace(".", "")
    
    # Image formats
    image_formats = ["png", "jpg", "jpeg", "webp", "bmp", "gif", "ico"]
    if ext in image_formats:
        return [f for f in image_formats if f != ext and f != ("jpeg" if ext == "jpg" else "jpg")]
        
    # Data formats
    if ext == "csv":
        return ["json", "txt", "xlsx"]
    if ext == "json":
        return ["csv", "txt"]
    if ext == "xlsx":
        return ["csv", "json"]

    # Document formats
    if ext == "pdf":
        return ["docx", "txt"]
    if ext == "docx":
        return ["pdf", "txt"]
    if ext == "txt":
        return ["csv", "json", "docx", "pdf"]
        
    # Presentation formats
    if ext == "pptx":
        return ["pdf", "txt"]
        
    return []

def launch_format_converter(root):
    t = config.themes[config.settings["theme_index"]]

    win = ctk.CTkToplevel(root)
    win.title("🔄 Format Converter")
    win.geometry("600x500")
    win.minsize(500, 450)
    win.configure(fg_color=t["bg"])
    win.attributes("-topmost", True)

    current_file_path: list[str | None] = [None]
    
    # ══════════════════════ HEADER ══════════════════════
    header = ctk.CTkFrame(win, fg_color=t["sidebar"], height=50, corner_radius=0)
    header.pack(fill="x")
    ctk.CTkLabel(header, text="🔄  Format Converter", font=("Segoe UI", 18, "bold"),
                 text_color=t["accent"]).pack(side="left", padx=20, pady=12)

    status_label = ctk.CTkLabel(header, text="Ready", font=("Segoe UI", 11),
                                text_color=t["muted"])
    status_label.pack(side="right", padx=20)

    # ══════════════════════ CONTENT ══════════════════════
    content = ctk.CTkFrame(win, fg_color="transparent")
    content.pack(fill="both", expand=True, padx=20, pady=20)

    # FILE SELECTION
    file_frame = ctk.CTkFrame(content, fg_color=t["card"], corner_radius=10)
    file_frame.pack(fill="x", pady=(0, 15))

    file_label = ctk.CTkLabel(file_frame, text="No file selected", font=("Segoe UI", 13), text_color=t["muted"], wraplength=400)
    file_label.pack(side="left", padx=15, pady=15, expand=True, anchor="w")

    def select_file():
        path = filedialog.askopenfilename(title="Select a file to convert")
        if path:
            current_file_path[0] = path
            file_label.configure(text=os.path.basename(path), text_color=t["fg"])
            _, ext = os.path.splitext(path)
            
            supported = list_supported_formats(ext)
            if not supported:
                status_label.configure(text="⚠️ Format not supported", text_color=t.get("danger", "#f87171"))
                target_dropdown.configure(values=["Unavailable"], state="disabled")
                target_var.set("Unavailable")
                convert_btn.configure(state="disabled")
            else:
                status_label.configure(text="Ready", text_color=t["muted"])
                target_dropdown.configure(values=[f.upper() for f in supported], state="normal")
                target_var.set(supported[0].upper())
                convert_btn.configure(state="normal")

    select_btn = ctk.CTkButton(file_frame, text="📂 Choose File", width=120, command=select_file,
                               fg_color=t["btn"], hover_color=t["btn_hover"], text_color=t["fg"])
    select_btn.pack(side="right", padx=15, pady=15)

    # TARGET FORMAT
    target_frame = ctk.CTkFrame(content, fg_color="transparent")
    target_frame.pack(fill="x", pady=10)
    
    ctk.CTkLabel(target_frame, text="Convert to:", font=("Segoe UI", 14, "bold"), text_color=t["fg"]).pack(side="left", padx=(0, 15))
    
    target_var = ctk.StringVar(value="---")
    target_dropdown = ctk.CTkOptionMenu(
        target_frame, variable=target_var, values=["---"], state="disabled",
        width=150, height=35, corner_radius=8,
        fg_color=t["btn"], button_color=t["accent"], button_hover_color=t["accent_hover"],
        text_color=t["fg"], font=("Segoe UI", 13, "bold"),
    )
    target_dropdown.pack(side="left")

    # PROGRESS & INFO
    info_frame = ctk.CTkFrame(content, fg_color="transparent")
    info_frame.pack(fill="both", expand=True, pady=10)
    
    info_label = ctk.CTkLabel(info_frame, text="Supported:\nImages (PNG, JPG, WEBP, GIF, BMP, ICO)\nData (CSV, JSON, XLSX)\nDocs (PDF, DOCX, TXT, PPTX)", 
                              font=("Segoe UI", 12), text_color=t["muted"], justify="center")
    info_label.pack(expand=True)

    # ══════════════════════ BOTTOM BAR ══════════════════════
    bottom = ctk.CTkFrame(win, fg_color=t["sidebar"], height=60, corner_radius=0)
    bottom.pack(fill="x", side="bottom")

    def perform_conversion():
        if not current_file_path[0]:
            return
            
        src_path = current_file_path[0]
        target_fmt = target_var.get().lower()
        if target_fmt == "---" or target_fmt == "unavailable" or not src_path:
            return
            
        _, src_ext = os.path.splitext(src_path)
        src_ext = src_ext.lower().replace(".", "")
        
        default_name = f"{os.path.splitext(os.path.basename(src_path))[0]}.{target_fmt}"
        save_path = filedialog.asksaveasfilename(
            title="Save Converted File",
            initialfile=default_name,
            defaultextension=f".{target_fmt}",
            filetypes=[(f"{target_fmt.upper()} File", f"*.{target_fmt}")]
        )
        
        if not save_path:
            return
            
        convert_btn.configure(state="disabled", text="⏳ Converting...")
        status_label.configure(text="Converting...", text_color=t["accent"])
        info_label.configure(text="Processing...", text_color=t["fg"])
        
        def _conv():
            try:
                # Images
                img_formats = ["png", "jpg", "jpeg", "webp", "bmp", "gif", "ico"]
                if src_ext in img_formats and target_fmt in img_formats:
                    img = Image.open(src_path)
                    if target_fmt in ["jpg", "jpeg"]:
                        img = img.convert("RGB")
                    img.save(save_path)
                    
                # Data
                elif src_ext == "csv" and target_fmt == "json":
                    pd.read_csv(src_path).to_json(save_path, orient="records", indent=4)
                elif src_ext == "csv" and target_fmt == "xlsx":
                    pd.read_csv(src_path).to_excel(save_path, index=False)
                elif src_ext == "json" and target_fmt == "csv":
                    pd.read_json(src_path).to_csv(save_path, index=False)
                elif src_ext == "xlsx" and target_fmt == "csv":
                    pd.read_excel(src_path).to_csv(save_path, index=False)
                elif src_ext == "xlsx" and target_fmt == "json":
                    pd.read_excel(src_path).to_json(save_path, orient="records", indent=4)
                
                elif src_ext == "pdf" and target_fmt == "docx":
                    cv = Converter(src_path)
                    cv.convert(save_path, start=0, end=None)
                    cv.close()
                elif src_ext == "docx" and target_fmt == "pdf":
                    convert(src_path, save_path)
                elif src_ext == "pdf" and target_fmt == "txt":
                    with open(src_path, "rb") as f, open(save_path, "w", encoding="utf-8") as out:
                        reader = PyPDF2.PdfReader(f)
                        for page in reader.pages:
                            out.write(page.extract_text() + "\n")
                elif src_ext == "docx" and target_fmt == "txt":
                    doc = docx.Document(src_path)
                    with open(save_path, "w", encoding="utf-8") as out:
                        for p in doc.paragraphs:
                            out.write(p.text + "\n")
                elif src_ext == "pptx" and target_fmt == "pdf":
                    if COMTYPES_AVAILABLE:
                        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                        powerpoint.Visible = 1
                        try:
                            deck = powerpoint.Presentations.Open(os.path.abspath(src_path))
                            deck.SaveAs(os.path.abspath(save_path), 32) # 32 is pdf
                            deck.Close()
                        finally:
                            powerpoint.Quit()
                    else:
                        raise ValueError("comtypes is required for pptx to pdf conversion on Windows.")
                elif src_ext == "pptx" and target_fmt == "txt":
                    prs = Presentation(src_path)
                    with open(save_path, "w", encoding="utf-8") as out:
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    out.write(str(shape.text) + "\n")
                elif target_fmt == "txt":
                    # Generic fallback
                    shutil.copyfile(src_path, save_path)
                else:
                    raise ValueError("Conversion not supported yet")
                    
                win.after(0, lambda: _finish("✅ Successfully converted!", t.get("success", "#22c55e")))
            except Exception as e:
                win.after(0, lambda: _finish(f"⚠️ Error: {e}", t.get("danger", "#f87171")))

        def _finish(msg, color):
            convert_btn.configure(state="normal", text="🔄 Convert & Save")
            status_label.configure(text=msg, text_color=color)
            info_label.configure(text=msg, text_color=color)
            win.after(4000, lambda: status_label.configure(text="Ready", text_color=t["muted"]))

        threading.Thread(target=_conv, daemon=True).start()

    convert_btn = ctk.CTkButton(
        bottom, text="🔄 Convert & Save", height=40, width=160,
        corner_radius=10, font=("Segoe UI", 14, "bold"),
        fg_color=t["accent"], hover_color=t["accent_hover"], text_color="#fff",
        state="disabled", command=perform_conversion
    )
    convert_btn.pack(pady=10)

def handle_prompt_conversion(src_path, prompt, callback):
    """
    Handle NLP based conversion requests directly from main.py
    """
    prompt = prompt.lower()
    
    # Simple NLP logic to extract target format
    target_formats = ["png", "jpg", "jpeg", "webp", "bmp", "gif", "ico", 
                      "csv", "json", "xlsx", "pdf", "docx", "txt", "pptx"]
                      
    target_fmt = None
    for fmt in target_formats:
        if f"to {fmt}" in prompt or f"in {fmt}" in prompt or f"into {fmt}" in prompt or fmt in prompt.replace(".", " ").split():
            target_fmt = fmt
            break
            
    if not target_fmt:
        callback("⚠️ Please specify a valid format to convert to (e.g. 'convert to pdf').", True)
        return
        
    if not src_path:
        callback("⚠️ Target file path not valid.", True)
        return
        
    _, src_ext = os.path.splitext(src_path)
    src_ext = src_ext.lower().replace(".", "")
    
    if target_fmt not in list_supported_formats(src_ext):
        callback(f"⚠️ Cannot convert from {src_ext.upper()} to {target_fmt.upper()}.", True)
        return
        
    callback("⏳ Converting file, please wait...", False)
    
    def _conv():
        try:
            download_dir = os.path.join(os.path.expanduser("~"), "Downloads")
            default_name = f"{os.path.splitext(os.path.basename(src_path))[0]}_converted.{target_fmt}"
            save_path = os.path.join(download_dir, default_name)
            
            # Images
            img_formats = ["png", "jpg", "jpeg", "webp", "bmp", "gif", "ico"]
            if src_ext in img_formats and target_fmt in img_formats:
                img = Image.open(src_path)
                if target_fmt in ["jpg", "jpeg"]:
                    img = img.convert("RGB")
                img.save(save_path)
                
            # Data
            elif src_ext == "csv" and target_fmt == "json":
                pd.read_csv(src_path).to_json(save_path, orient="records", indent=4)
            elif src_ext == "csv" and target_fmt == "xlsx":
                pd.read_csv(src_path).to_excel(save_path, index=False)
            elif src_ext == "json" and target_fmt == "csv":
                pd.read_json(src_path).to_csv(save_path, index=False)
            elif src_ext == "xlsx" and target_fmt == "csv":
                pd.read_excel(src_path).to_csv(save_path, index=False)
            elif src_ext == "xlsx" and target_fmt == "json":
                pd.read_excel(src_path).to_json(save_path, orient="records", indent=4)
            
            elif src_ext == "pdf" and target_fmt == "docx":
                cv = Converter(src_path)
                cv.convert(save_path, start=0, end=None)
                cv.close()
            elif src_ext == "docx" and target_fmt == "pdf":
                convert(src_path, save_path)
            elif src_ext == "pdf" and target_fmt == "txt":
                with open(src_path, "rb") as f, open(save_path, "w", encoding="utf-8") as out:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        out.write(page.extract_text() + "\\n")
            elif src_ext == "docx" and target_fmt == "txt":
                doc = docx.Document(src_path)
                with open(save_path, "w", encoding="utf-8") as out:
                    for p in doc.paragraphs:
                        out.write(p.text + "\\n")
            elif src_ext == "pptx" and target_fmt == "pdf":
                if COMTYPES_AVAILABLE:
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    powerpoint.Visible = 1
                    try:
                        deck = powerpoint.Presentations.Open(os.path.abspath(src_path))
                        deck.SaveAs(os.path.abspath(save_path), 32)
                        deck.Close()
                    finally:
                        powerpoint.Quit()
                else:
                    raise ValueError("comtypes is required for pptx to pdf conversion on Windows.")
            elif src_ext == "pptx" and target_fmt == "txt":
                prs = Presentation(src_path)
                with open(save_path, "w", encoding="utf-8") as out:
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                out.write(str(shape.text) + "\\n")
            elif target_fmt == "txt":
                shutil.copyfile(src_path, save_path)
            else:
                raise ValueError("Conversion logic missing")
                
            callback(f"✅ Successfully converted your file!\\n\\nSaved to: `{save_path}`", True)
        except Exception as e:
            callback(f"⚠️ Failed to convert file: {e}", True)

    threading.Thread(target=_conv, daemon=True).start()
